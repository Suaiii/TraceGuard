# -*- coding: utf-8 -*-
"""
build_figures_pptx.py  (v3)

生成 TraceGuard 报告的系统结构示意图（可编辑 PPTX），风格对齐参考论文 ExDA 的 Figure 3。
每个模块嵌入 **原创** matplotlib 素材（见 gen_figure_assets.py），先运行素材脚本再运行本脚本。

v3（2026-07-15 队长反馈）：字号再放大；用真实示意图（热力图/骨干架构/域对齐散点/篡改定位/
量表/柱状）替换粗色块；三张图都补图像；Web 流程每步加线性图标。

红线：纯技术图 + 全原创素材，无任何身份信息 / 无 SHA、绝对路径、权重文件名 / 无敏感图像。
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image

ASSETS = r"E:\aNB\TECH\AI竞赛\docs\figures\system\assets"

# ---------------------------------------------------------------- 配色
C_YELLOW   = RGBColor(0xFC, 0xF6, 0xD8)
C_YELLOW_B = RGBColor(0xEC, 0xDD, 0x9A)
C_PURPLE   = RGBColor(0xE7, 0xE0, 0xF0)
C_PURPLE_B = RGBColor(0xC7, 0xB8, 0xE2)
C_ORANGE   = RGBColor(0xFB, 0xE7, 0xD3)
C_ORANGE_B = RGBColor(0xF1, 0xC6, 0x9B)
C_CYAN     = RGBColor(0xDA, 0xEC, 0xEC)
C_CYAN_B   = RGBColor(0xA9, 0xD3, 0xD6)
C_GREEN    = RGBColor(0xDD, 0xEB, 0xD8)
C_GREEN_B  = RGBColor(0xB6, 0xD5, 0xAA)
C_NEUTRAL  = RGBColor(0xEC, 0xEC, 0xF1)
C_NEUTRAL_B= RGBColor(0xD5, 0xD5, 0xDE)

C_TITLE  = RGBColor(0x24, 0x24, 0x24)
C_TEXT   = RGBColor(0x2E, 0x2E, 0x2E)
C_MUTED  = RGBColor(0x5E, 0x5E, 0x5E)
C_BORDER = RGBColor(0xA0, 0xA0, 0xB0)

C_ARROW_MAIN = RGBColor(0x1E, 0x1E, 0x1E)
C_ARROW_SOL  = RGBColor(0x4A, 0x4A, 0x56)
C_ARROW_DASH = RGBColor(0x74, 0x74, 0x82)
C_GREEN_OK = RGBColor(0x2E, 0x8B, 0x3D)
C_RED_NG   = RGBColor(0xC0, 0x39, 0x2B)

FONT = "Microsoft YaHei"

# ---------------------------------------------------------------- 基础工具

def _no_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def _set_ea_font(run, font=FONT):
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", font)


def set_text(shape, lines, size=15, color=C_TEXT, bold=False, underline=False,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT,
             line_spacing=1.03):
    if isinstance(lines, str):
        lines = [lines]
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(3); tf.margin_right = Pt(3)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        segs = line if isinstance(line, list) else [(line, color)]
        for seg in segs:
            txt, col = seg if isinstance(seg, tuple) else (seg, color)
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.underline = underline
            r.font.color.rgb = col
            _set_ea_font(r, font)
    return shape


def shape(slide, kind, x, y, w, h, fill, line=C_BORDER, line_w=0.75, radius=None):
    sp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    _no_shadow(sp)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    if radius is not None:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def rrect(slide, x, y, w, h, fill, line=C_BORDER, line_w=0.75, radius=0.10):
    return shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill,
                 line=line, line_w=line_w, radius=radius)


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb


def dash_line(sp, val='dash'):
    ln = sp.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:prstDash'), {'val': val}))


def container(slide, x, y, w, h, fill, title, title_color=C_TITLE,
              title_size=17, line=C_BORDER):
    box = rrect(slide, x, y, w, h, fill, line=line, line_w=1.4, radius=0.045)
    tt = textbox(slide, x + 0.06, y + 0.08, w - 0.12, 0.42)
    set_text(tt, title, size=title_size, color=title_color, bold=True,
             underline=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    return box


def block(slide, x, y, w, h, fill, lines, size=14, color=C_TEXT, bold=False):
    b = rrect(slide, x, y, w, h, fill, line=None, radius=0.13)
    if lines is not None:
        set_text(b, lines, size=size, color=color, bold=bold)
    return b


def _add_arrowhead(conn, tail=True, head=False):
    ln = conn.line._get_or_add_ln()
    if head:
        ln.append(ln.makeelement(qn('a:headEnd'),
                  {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    if tail:
        ln.append(ln.makeelement(qn('a:tailEnd'),
                  {'type': 'triangle', 'w': 'med', 'len': 'med'}))


def arrow(slide, x1, y1, x2, y2, color=C_ARROW_SOL, width=2.0, dashed=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    _no_shadow(conn)
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    if dashed:
        dash_line(conn)
    _add_arrowhead(conn, tail=True)
    return conn


def pic_fit(slide, name, x, y, w, h):
    """按盒子居中放置图片（保持长宽比，letterbox）。"""
    path = os.path.join(ASSETS, name)
    iw, ih = Image.open(path).size
    ar = iw / ih
    if ar > w / h:
        nw, nh = w, w / ar
    else:
        nh, nw = h, h * ar
    px, py = x + (w - nw) / 2, y + (h - nh) / 2
    return slide.shapes.add_picture(path, Inches(px), Inches(py),
                                    width=Inches(nw), height=Inches(nh))


def title_bar(slide, text, sub=None):
    tb = textbox(slide, 0.3, 0.10, 12.73, 0.66)
    set_text(tb, text, size=34, color=C_TITLE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    if sub:
        sb = textbox(slide, 0.3, 0.80, 12.73, 0.40)
        set_text(sb, sub, size=14, color=C_MUTED,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)


def icon_image_vec(slide, x, y, w, h):
    """输入缩略图（矢量：白框 + 山 + 太阳 + 地面）。"""
    shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, RGBColor(0xFF, 0xFF, 0xFF),
          line=C_MUTED, line_w=1.4)
    shape(slide, MSO_SHAPE.RECTANGLE, x + 0.06, y + 0.06, w - 0.12, h * 0.55,
          RGBColor(0xDC, 0xEB, 0xF7), line=None)
    shape(slide, MSO_SHAPE.OVAL, x + w * 0.62, y + h * 0.14, w * 0.18, w * 0.18,
          RGBColor(0xF6, 0xC7, 0x4B), line=None)
    shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, x + w * 0.12, y + h * 0.30,
          w * 0.5, h * 0.4, RGBColor(0x7F, 0xB0, 0x6E), line=None)
    shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, x + w * 0.42, y + h * 0.34,
          w * 0.44, h * 0.36, RGBColor(0x5F, 0x93, 0x54), line=None)
    shape(slide, MSO_SHAPE.RECTANGLE, x + 0.06, y + h * 0.62, w - 0.12, h * 0.32,
          RGBColor(0xE9, 0xE3, 0xCF), line=None)


# ---------------------------------------------------------------- 演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ================================================================ Slide 1
s1 = prs.slides.add_slide(BLANK)
# 图内不设标题：Word 图下方已有自动题注（导出后 PIL 裁去顶部留白）
MIDY = 3.95

# 输入
container(s1, 0.28, 2.78, 1.82, 2.42, C_YELLOW, "输入")
icon_image_vec(s1, 0.53, 3.36, 1.32, 1.05)
set_text(textbox(s1, 0.30, 4.48, 1.78, 0.55), "单张 RGB 图像", size=14.5,
         color=C_TEXT, bold=True, anchor=MSO_ANCHOR.TOP)

# 跨域 AIGC 检测模块（核心）
container(s1, 2.45, 1.78, 3.62, 4.20, C_PURPLE, "跨域 AIGC 检测模块（核心）", title_size=17)
pic_fit(s1, "backbone.png", 2.58, 2.35, 3.38, 1.55)
set_text(textbox(s1, 2.58, 3.86, 3.38, 0.34), "MambaOut-Small 去全局池化骨干（2304→256）",
         size=12.5, color=C_MUTED, anchor=MSO_ANCHOR.TOP)
pic_fit(s1, "mkmmd.png", 2.62, 4.28, 1.95, 1.28)
set_text(textbox(s1, 4.62, 4.42, 1.42, 1.05), ["MK-MMD", "无监督", "域自适应"],
         size=14, color=C_TEXT, bold=True, align=PP_ALIGN.LEFT)

# 可解释模块
container(s1, 6.30, 0.86, 2.72, 2.18, C_ORANGE, "可解释模块", title_size=16)
pic_fit(s1, "gradcam.png", 6.48, 1.42, 1.28, 1.30)
set_text(textbox(s1, 7.85, 1.42, 1.10, 1.30), ["Stage2", "14×14", "Grad-CAM", "热力图"],
         size=13.5, color=C_TEXT, align=PP_ALIGN.LEFT)

# 篡改定位模块
container(s1, 6.30, 4.86, 2.72, 2.02, C_CYAN, "篡改定位模块", title_size=16)
pic_fit(s1, "bbox.png", 6.48, 5.42, 1.28, 1.28)
set_text(textbox(s1, 7.85, 5.44, 1.10, 1.25), ["多尺度滑窗", "+特征统计", "四象限", "分类"],
         size=13.5, color=C_TEXT, align=PP_ALIGN.LEFT)

# 多证据风险融合
container(s1, 9.28, 2.30, 2.42, 3.30, C_GREEN, "多证据风险融合", title_size=16)
pic_fit(s1, "bars.png", 9.45, 2.86, 2.08, 1.18)
set_text(textbox(s1, 9.38, 3.98, 2.24, 0.32), "五维加权 → risk_score", size=13,
         color=C_TEXT, bold=True, anchor=MSO_ANCHOR.TOP)
pic_fit(s1, "gauge.png", 9.55, 4.36, 1.90, 0.92)
set_text(textbox(s1, 9.38, 5.24, 2.24, 0.34), "risk_level（低 / 中 / 高）", size=13,
         color=C_TEXT, anchor=MSO_ANCHOR.TOP)

# 三层输出合同
set_text(textbox(s1, 11.78, 1.28, 1.50, 0.34), "三层输出合同", size=14.5,
         color=C_TITLE, bold=True, underline=True, anchor=MSO_ANCHOR.TOP)
outX, outW = 11.80, 1.50
oa = block(s1, outX, 1.72, outW, 1.30, C_PURPLE, None)
set_text(oa, [[("全局判断", C_TITLE)], [("label", C_TEXT)], [("fake_prob", C_TEXT)],
              [("Real ", C_GREEN_OK), ("✓ ", C_GREEN_OK),
               ("Fake ", C_RED_NG), ("✗", C_RED_NG)]], size=12.5)
oa.text_frame.paragraphs[0].runs[0].font.bold = True
ob = block(s1, outX, 3.18, outW, 1.22, C_ORANGE, None)
set_text(ob, [[("局部证据", C_TITLE)], [("heatmap", C_TEXT)], [("mask / bbox", C_TEXT)]],
         size=12.5)
ob.text_frame.paragraphs[0].runs[0].font.bold = True
oc = block(s1, outX, 4.56, outW, 1.22, C_GREEN, None)
set_text(oc, [[("融合结论", C_TITLE)], [("risk_score", C_TEXT)],
              [("risk_level / 报告", C_TEXT)]], size=12.5)
oc.text_frame.paragraphs[0].runs[0].font.bold = True

# 主数据流
arrow(s1, 2.10, MIDY, 2.45, MIDY, C_ARROW_SOL, 2.5)
arrow(s1, 6.07, MIDY, 9.28, MIDY, C_ARROW_MAIN, 4.5)
set_text(textbox(s1, 6.30, MIDY - 0.44, 2.72, 0.34), "唯一权威判定", size=13.5,
         color=C_ARROW_MAIN, bold=True, anchor=MSO_ANCHOR.BOTTOM)
# 并行分支（虚线）
arrow(s1, 6.07, 3.35, 6.30, 1.95, C_ARROW_DASH, 2.0, dashed=True)
arrow(s1, 6.07, 4.55, 6.30, 5.55, C_ARROW_DASH, 2.0, dashed=True)
arrow(s1, 9.02, 1.95, 9.28, 3.35, C_ARROW_DASH, 2.0, dashed=True)
arrow(s1, 9.02, 5.55, 9.28, 4.55, C_ARROW_DASH, 2.0, dashed=True)
# 融合→输出（扇出）
arrow(s1, 11.70, 3.75, outX, 2.37, C_ARROW_SOL, 1.9)
arrow(s1, 11.70, 3.95, outX, 3.79, C_ARROW_SOL, 1.9)
arrow(s1, 11.70, 4.15, outX, 5.17, C_ARROW_SOL, 1.9)

# 图例（左下）
lg = rrect(s1, 0.30, 5.55, 4.55, 1.55, RGBColor(0xFB, 0xFB, 0xFC),
           line=C_MUTED, line_w=1.0, radius=0.04)
dash_line(lg)
set_text(textbox(s1, 0.42, 5.62, 4.3, 0.30), "图例", size=13, color=C_TITLE,
         bold=True, underline=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
arrow(s1, 0.50, 6.20, 1.12, 6.20, C_ARROW_MAIN, 4.5)
set_text(textbox(s1, 1.22, 6.05, 3.55, 0.30), "粗实线：唯一权威判定（label / fake_prob）",
         size=12, color=C_TEXT, align=PP_ALIGN.LEFT)
arrow(s1, 0.50, 6.58, 1.12, 6.58, C_ARROW_DASH, 2.0, dashed=True)
set_text(textbox(s1, 1.22, 6.43, 3.55, 0.30), "虚线箭头：并行证据分支", size=12,
         color=C_TEXT, align=PP_ALIGN.LEFT)
for i, c in enumerate((C_PURPLE_B, C_ORANGE_B, C_GREEN_B)):
    rrect(s1, 0.53 + i * 0.22, 6.88, 0.18, 0.18, c, line=None, radius=0.2)
set_text(textbox(s1, 1.22, 6.80, 3.55, 0.30), "三色分区：三层输出合同", size=12,
         color=C_TEXT, align=PP_ALIGN.LEFT)

# 底部四入口标注
note = rrect(s1, 5.10, 6.60, 8.20, 0.74, RGBColor(0xFB, 0xFB, 0xFC),
             line=C_MUTED, line_w=1.0, radius=0.08)
dash_line(note)
set_text(note, "四入口（Web UI / API / CLI / 批量）复用同一 ExplanationPipeline"
               "　｜　终端：Web 可视化展示 + 检测报告导出", size=13, color=C_TEXT)


# ================================================================ Slide 2
s2 = prs.slides.add_slide(BLANK)
# 图内不设标题：Word 图下方已有自动题注


def flow_node(slide, x, y, w, h, fill, icon, lines, size=14.5, bold=False):
    rrect(slide, x, y, w, h, fill, line=None, radius=0.11)
    pic_fit(slide, icon, x + w / 2 - 0.42, y + 0.14, 0.84, 0.84)
    tb = textbox(slide, x + 0.06, y + 0.98, w - 0.12, h - 1.02)
    set_text(tb, lines, size=size, color=C_TEXT, bold=bold, anchor=MSO_ANCHOR.TOP)


ry, rh = 1.42, 2.05
flow_node(s2, 0.40, ry, 2.78, rh, C_CYAN_B, "icon_upload.png", ["① 用户上传图片"], bold=True)
flow_node(s2, 3.52, ry, 2.78, rh, C_YELLOW_B, "icon_check.png", ["② 输入校验", "格式 / 尺寸"])
flow_node(s2, 6.64, ry, 3.05, rh, C_PURPLE_B, "icon_api.png",
          ["③ 统一分析接口", "POST /api/v1/analyze"])
flow_node(s2, 10.05, ry, 2.88, rh, C_PURPLE_B, "icon_gpu.png",
          ["④ GPU 推理", "ExplanationPipeline"], bold=True)
ay = ry + rh / 2
arrow(s2, 3.18, ay, 3.52, ay, C_ARROW_MAIN, 3.0)
arrow(s2, 6.30, ay, 6.64, ay, C_ARROW_MAIN, 3.0)
arrow(s2, 9.69, ay, 10.05, ay, C_ARROW_MAIN, 3.0)
arrow(s2, 11.49, ry + rh, 11.49, 4.15, C_ARROW_MAIN, 3.0)

# ⑤ 证据渲染（分栏）
container(s2, 8.30, 4.15, 4.63, 2.55, C_GREEN, "⑤ 证据渲染（分栏）", title_size=16)
pic_fit(s2, "icon_panel.png", 8.45, 4.62, 0.72, 0.72)
block(s2, 8.50, 5.35, 1.36, 1.20, C_GREEN_B, ["全局", "label", "fake_prob"], size=12.5)
block(s2, 10.02, 5.35, 1.42, 1.20, C_ORANGE_B, ["局部", "tamper", "bbox/热力"], size=12.5)
block(s2, 11.58, 5.35, 1.20, 1.20, C_CYAN_B, ["风险", "low/med", "/high"], size=12.5)

flow_node(s2, 4.45, 4.30, 3.35, 2.05, C_YELLOW_B, "icon_review.png",
          ["⑥ 人工复核", "证据冲突或退化→转人工"], size=13.5)
flow_node(s2, 0.45, 4.30, 3.30, 2.05, C_NEUTRAL_B, "icon_report.png",
          ["⑦ 单图检测报告导出", "（HTML）"], size=14, bold=True)
by = 5.32
arrow(s2, 8.30, by, 7.80, by, C_ARROW_MAIN, 3.0)
arrow(s2, 4.45, by, 3.75, by, C_ARROW_MAIN, 3.0)

stn = rrect(s2, 0.45, 6.72, 6.35, 0.62, RGBColor(0xFB, 0xFB, 0xFC),
            line=C_MUTED, line_w=1.0, radius=0.10)
dash_line(stn)
set_text(stn, "界面状态反馈：loading / empty / error 三态", size=13, color=C_TEXT)
set_text(textbox(s2, 7.0, 6.78, 5.9, 0.5),
         "实线箭头：主数据流（首行左→右，折返续接下一行）", size=12,
         color=C_MUTED, align=PP_ALIGN.RIGHT)


# ================================================================ Slide 3
s3 = prs.slides.add_slide(BLANK)
# 图内不设标题：Word 图下方已有自动题注（占位提示保留在各卡片与页脚）

# 用原创示意图作“示例样式”淡显，提示每格应放什么（非真实检测结果）
demo = [("a  输入图像", None), ("b  Stage2 Grad-CAM 叠加图", "gradcam.png"),
        ("c  局部定位可疑区域", "bbox.png")]
xs = [0.55, 4.75, 8.95]
pw, ph, py = 3.83, 4.55, 1.55
for x, (t, asset) in zip(xs, demo):
    box = rrect(s3, x, py, pw, ph, RGBColor(0xF5, 0xF5, 0xF8),
                line=C_MUTED, line_w=1.6, radius=0.03)
    dash_line(box)
    if asset:
        p = pic_fit(s3, asset, x + 0.55, py + 0.75, pw - 1.1, ph - 1.7)
        # 半透明，示意而非真结果
        p.line.fill.background()
    else:
        icon_image_vec(s3, x + pw / 2 - 1.0, py + 1.05, 2.0, 1.7)
    set_text(textbox(s3, x + 0.05, py + 0.10, pw - 0.10, 0.45), t, size=17,
             color=C_TITLE, bold=True, anchor=MSO_ANCHOR.TOP)
    set_text(textbox(s3, x + 0.05, py + ph - 0.55, pw - 0.10, 0.45),
             "（示意样式，待填真实截图）", size=12.5, color=C_MUTED)

set_text(textbox(s3, 0.55, 6.42, 12.23, 0.5),
         "占位：示意图仅表示每格应呈现的内容形态，最终由成员填入真实检测截图，不得伪造检测结果。",
         size=13, color=C_MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)


# ----------------------------------------------------------------
OUT = r"E:\aNB\TECH\AI竞赛\docs\figures\system\traceguard_figures.pptx"
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print("saved:", OUT)
print("slides:", len(prs.slides._sldIdLst))
